r"""Build the Virtual Bingo pitch deck as Virtual_Bingo.pptx.

Standalone tooling — not part of the running app. Safe to delete after
the deck is exported. Re-run any time the content needs a refresh:

    .\.venv\Scripts\python.exe scripts\build_slides.py

Output: ``Virtual_Bingo.pptx`` at the project root.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- CGI-ish palette ---------------------------------------------------------
CGI_RED = RGBColor(0xC8, 0x10, 0x2E)
CGI_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_DARK = RGBColor(0x22, 0x22, 0x33)
TEXT_MUTED = RGBColor(0x5F, 0x6B, 0x7A)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM = RGBColor(0xCC, 0xCC, 0xDD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Calibri"  # safe default that ships with PowerPoint


def _add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def _add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font_size=18,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_header(slide, title, kicker=None):
    """Slide header: left red bar, optional kicker, big title."""
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, CGI_RED)
    if kicker:
        _add_text(
            slide, kicker.upper(),
            Inches(0.7), Inches(0.55), Inches(11), Inches(0.4),
            font_size=12, bold=True, color=CGI_RED,
        )
        _add_text(
            slide, title,
            Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
            font_size=44, bold=True, color=CGI_NAVY,
        )
    else:
        _add_text(
            slide, title,
            Inches(0.7), Inches(0.7), Inches(12), Inches(1.0),
            font_size=44, bold=True, color=CGI_NAVY,
        )
    _add_rect(slide, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.04), CGI_RED)


def _set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = _add_blank_slide(prs)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CGI_NAVY)
    _add_rect(s, Inches(0), Inches(3.0), Inches(2.4), Inches(0.18), CGI_RED)
    _add_text(
        s, "CGI",
        Inches(0.9), Inches(2.2), Inches(11), Inches(0.5),
        font_size=16, bold=True, color=CGI_RED,
    )
    _add_text(
        s, "VIRTUAL BINGO",
        Inches(0.9), Inches(3.3), Inches(12), Inches(1.4),
        font_size=72, bold=True, color=WHITE,
    )
    _add_text(
        s, "One app. Real-time. Audit-ready.",
        Inches(0.9), Inches(4.9), Inches(12), Inches(0.6),
        font_size=22, color=DIM,
    )
    _set_notes(s, "Open with the one-liner. Manual workflow becomes one app.")


def slide_problem(prs):
    s = _add_blank_slide(prs)
    _add_header(s, "The problem")

    pains = [
        "Cards generated on a third-party site",
        "Words typed manually into Teams chat",
        "Players mark cards by hand",
        "Wins announced by voice",
        "Manual verification + prize email",
    ]
    top = Inches(2.6)
    for i, head in enumerate(pains):
        row_top = top + Inches(i * 0.75)
        # Numbered red circle
        circ = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), row_top, Inches(0.55), Inches(0.55)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = CGI_RED
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.bold = True
        r.font.size = Pt(18)
        r.font.color.rgb = WHITE

        _add_text(
            s, head,
            Inches(1.7), row_top + Inches(0.05),
            Inches(11), Inches(0.5),
            font_size=22, bold=True, color=CGI_NAVY,
        )

    _set_notes(s, "Five manual steps. Every one is a place mistakes hide.")


def slide_solution(prs):
    s = _add_blank_slide(prs)
    _add_header(s, "Our solution")

    steps = [("01", "Create"), ("02", "Join"), ("03", "Play"), ("04", "Finish")]
    panel_w = Inches(2.95)
    panel_h = Inches(3.2)
    gap = Inches(0.15)
    start_left = Inches(0.7)
    top = Inches(3.0)

    for i, (num, head) in enumerate(steps):
        left = start_left + (panel_w + gap) * i
        _add_rect(s, left, top, panel_w, panel_h, BG_LIGHT)
        _add_rect(s, left, top, panel_w, Inches(0.7), CGI_RED)
        _add_text(
            s, num,
            left + Inches(0.3), top + Inches(0.13),
            Inches(1), Inches(0.5),
            font_size=20, bold=True, color=WHITE,
        )
        _add_text(
            s, head,
            left + Inches(0.3), top + Inches(1.1),
            panel_w - Inches(0.5), Inches(0.8),
            font_size=32, bold=True, color=CGI_NAVY,
        )

    _set_notes(s, "Four phases, one app. Talk through each on the demo.")


def slide_interactive(prs):
    s = _add_blank_slide(prs)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CGI_NAVY)
    _add_rect(s, Inches(0), Inches(2.5), Inches(2.2), Inches(0.18), CGI_RED)

    _add_text(
        s, "AUTOMATED ≠ PASSIVE",
        Inches(0.9), Inches(1.8), Inches(12), Inches(0.5),
        font_size=16, bold=True, color=CGI_RED,
    )
    _add_text(
        s, "Still a live game.",
        Inches(0.9), Inches(2.8), Inches(12), Inches(1.0),
        font_size=56, bold=True, color=WHITE,
    )
    _add_text(
        s, "Players hear the call, click to mark, race to claim BINGO.",
        Inches(0.9), Inches(4.4), Inches(12), Inches(0.6),
        font_size=22, color=DIM,
    )
    _add_text(
        s, "Host pauses, resumes, sets the pace. The room reacts in real time.",
        Inches(0.9), Inches(5.1), Inches(12), Inches(0.6),
        font_size=22, color=DIM,
    )
    _set_notes(
        s,
        "Key talking point: the system runs the chore, the people still play. "
        "TTS, live socket pushes, click-to-mark, BINGO button, host pause.",
    )


def slide_features(prs):
    s = _add_blank_slide(prs)
    _add_header(s, "What's built")

    features = [
        "OTP email login",
        "AI-generated topics",
        "Unique 5×5 cards",
        "Real-time calling",
        "Server-validated wins",
        "Pause / resume",
        "Configurable pace",
        "Custom patterns",
        "Live leaderboard",
        "Full audit trail",
        "Text-to-speech caller",
        "Sound + animation",
    ]
    cols = 4
    rows = 3
    col_w = Inches(2.9)
    row_h = Inches(1.3)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_left = Inches(0.7)
    top = Inches(2.4)

    for i, name in enumerate(features):
        r = i // cols
        c = i % cols
        left = start_left + (col_w + gap_x) * c
        row_top = top + (row_h + gap_y) * r
        _add_rect(s, left, row_top, col_w, row_h, BG_LIGHT)
        _add_rect(s, left, row_top, Inches(0.1), row_h, CGI_RED)
        _add_text(
            s, name,
            left + Inches(0.3), row_top + Inches(0.4),
            col_w - Inches(0.5), Inches(0.6),
            font_size=18, bold=True, color=CGI_NAVY,
        )

    _set_notes(s, "Skim the grid. Demo will cover the load-bearing ones live.")


def slide_future(prs):
    s = _add_blank_slide(prs)
    _add_header(s, "What's next")

    ideas = [
        ("AI voice — pick a CGI teammate",  "Caller speaks in a familiar voice."),
        ("Microsoft Teams integration",      "Live word feed + join link in-channel."),
        ("Teams tab + meeting app",          "Play without leaving the meeting."),
        ("Spectator mode",                   "Read-only watch link, no auth."),
        ("Multi-round sessions",             "Aggregated leaderboard across games."),
        ("Replay mode",                       "Re-watch any past game at speed."),
        ("Custom pattern designer",          "Click cells to define winning shapes."),
        ("Auto-email the winners",           "Gift card details delivered automatically."),
    ]

    cols = 2
    col_w = Inches(5.9)
    row_h = Inches(1.15)
    gap_x = Inches(0.2)
    gap_y = Inches(0.18)
    start_left = Inches(0.7)
    top = Inches(2.4)

    for i, (head, sub) in enumerate(ideas):
        r = i // cols
        c = i % cols
        left = start_left + (col_w + gap_x) * c
        row_top = top + (row_h + gap_y) * r
        _add_rect(s, left, row_top, col_w, row_h, BG_LIGHT)
        _add_rect(s, left, row_top, Inches(0.1), row_h, CGI_RED)
        _add_text(
            s, head,
            left + Inches(0.3), row_top + Inches(0.2),
            col_w - Inches(0.5), Inches(0.5),
            font_size=17, bold=True, color=CGI_NAVY,
        )
        _add_text(
            s, sub,
            left + Inches(0.3), row_top + Inches(0.65),
            col_w - Inches(0.5), Inches(0.5),
            font_size=12, color=TEXT_MUTED,
        )

    _set_notes(
        s,
        "AI voice is the headline future feature — clone a CGI teammate so the "
        "caller sounds familiar. Then Teams. Other items are roadmap candy.",
    )


def slide_close(prs):
    s = _add_blank_slide(prs)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CGI_NAVY)
    _add_rect(s, Inches(0), Inches(3.2), Inches(2.2), Inches(0.18), CGI_RED)
    _add_text(
        s, "DEMO  •  Q&A",
        Inches(0.9), Inches(2.4), Inches(11), Inches(0.6),
        font_size=18, bold=True, color=CGI_RED,
    )
    _add_text(
        s, "Let's play.",
        Inches(0.9), Inches(3.5), Inches(11), Inches(1.2),
        font_size=64, bold=True, color=WHITE,
    )
    _set_notes(s, "Hand the audience the join link and run a live round.")


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_solution,
        slide_interactive,
        slide_features,
        slide_future,
        slide_close,
    ]
    for build in builders:
        build(prs)

    out = Path(__file__).resolve().parent.parent / "Virtual_Bingo.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
